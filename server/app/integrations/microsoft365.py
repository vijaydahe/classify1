"""Server-side Microsoft 365 auto-stamper.

Uses a Microsoft Entra (Azure AD) app registration with application permissions
(admin-consented once) to scan OneDrive/SharePoint via Microsoft Graph, classify
each Word document against the tenant's rules, and stamp the classification into
its footer — automatically, no user action. Mirrors the Google Drive integration.

Word .docx files are stamped by downloading, editing the footer with python-docx,
and uploading back through Graph.
"""
from __future__ import annotations

import io
import json
import urllib.error
import urllib.parse
import urllib.request
from datetime import datetime, timezone

from sqlalchemy.orm import Session

from .. import models
from ..classification import classify_text

GRAPH = "https://graph.microsoft.com/v1.0"
MAX_DOCS_PER_SCAN = 200


def _http(url: str, method: str = "GET", token: str | None = None,
          data: bytes | None = None, headers: dict | None = None, raw: bool = False):
    req = urllib.request.Request(url, data=data, method=method)
    if token:
        req.add_header("Authorization", "Bearer " + token)
    for k, v in (headers or {}).items():
        req.add_header(k, v)
    with urllib.request.urlopen(req, timeout=45) as resp:
        body = resp.read()
        if raw:
            return body
        return json.loads(body) if body else {}


def get_token(azure_tenant: str, client_id: str, client_secret: str) -> str:
    """App-only (client credentials) token for Microsoft Graph."""
    url = f"https://login.microsoftonline.com/{azure_tenant}/oauth2/v2.0/token"
    data = urllib.parse.urlencode({
        "grant_type": "client_credentials",
        "client_id": client_id,
        "client_secret": client_secret,
        "scope": "https://graph.microsoft.com/.default",
    }).encode()
    resp = _http(url, "POST", data=data,
                 headers={"Content-Type": "application/x-www-form-urlencoded"})
    return resp["access_token"]


def list_office_docs(token: str, drive_user: str) -> list[dict]:
    """Lists Word/Excel/PowerPoint files in the target user's OneDrive."""
    found = []
    for ext in (".docx", ".xlsx", ".pptx"):
        url = (f"{GRAPH}/users/{urllib.parse.quote(drive_user)}/drive/root/search(q='{ext}')"
               f"?$top={MAX_DOCS_PER_SCAN}&$select=id,name,file")
        for i in _http(url, token=token).get("value", []):
            if i.get("name", "").lower().endswith(ext):
                found.append(i)
    return found


def stamp_docx_bytes(data: bytes, text: str, color_hex: str, placement: str) -> bytes:
    """Stamps a .docx (in memory) into its header or footer, returns new bytes."""
    from docx import Document
    from docx.shared import Pt, RGBColor

    rgb = _rgb(color_hex)
    doc = Document(io.BytesIO(data))
    for section in doc.sections:
        container = section.header if placement == "header" else section.footer
        container.is_linked_to_previous = False
        para = container.paragraphs[0] if container.paragraphs else container.add_paragraph()
        if para.text.strip().startswith("CLASSIFICATION"):
            para.clear()
        run = para.add_run(text)
        run.font.bold = True
        run.font.size = Pt(10)
        run.font.color.rgb = rgb
    out = io.BytesIO()
    doc.save(out)
    return out.getvalue()


def stamp_xlsx_bytes(data: bytes, text: str, color_hex: str, placement: str) -> bytes:
    """Stamps an .xlsx into the print header/footer of every sheet (idempotent)."""
    from openpyxl import load_workbook
    wb = load_workbook(io.BytesIO(data))
    marker = "&\"-,Bold\"&K" + color_hex.lstrip("#").upper() + text
    for ws in wb.worksheets:
        target = ws.oddHeader.center if placement == "header" else ws.oddFooter.center
        target.text = marker
    out = io.BytesIO()
    wb.save(out)
    return out.getvalue()


def stamp_pptx_bytes(data: bytes, text: str, color_hex: str, placement: str) -> bytes:
    """Stamps a .pptx with a banner text box on every slide (idempotent)."""
    from pptx import Presentation
    from pptx.util import Emu, Pt
    from pptx.dml.color import RGBColor

    prs = Presentation(io.BytesIO(data))
    h = color_hex.lstrip("#")
    rgb = RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))
    for slide in prs.slides:
        if any(getattr(s, "has_text_frame", False) and "CLASSIFICATION" in s.text_frame.text
               for s in slide.shapes):
            continue
        top = Emu(120000) if placement == "header" else prs.slide_height - Emu(420000)
        box = slide.shapes.add_textbox(Emu(120000), top, prs.slide_width - Emu(240000), Emu(360000))
        run = box.text_frame.paragraphs[0].add_run()
        run.text = text
        run.font.bold = True
        run.font.size = Pt(10)
        run.font.color.rgb = rgb
    out = io.BytesIO()
    prs.save(out)
    return out.getvalue()


def stamp_office(data: bytes, name: str, text: str, color_hex: str, placement: str) -> bytes:
    n = name.lower()
    if n.endswith(".docx"):
        return stamp_docx_bytes(data, text, color_hex, placement)
    if n.endswith(".xlsx"):
        return stamp_xlsx_bytes(data, text, color_hex, placement)
    if n.endswith(".pptx"):
        return stamp_pptx_bytes(data, text, color_hex, placement)
    raise ValueError("unsupported type")


def _rgb(color_hex: str):
    from docx.shared import RGBColor
    h = color_hex.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def doc_text(data: bytes, name: str) -> str:
    """Extracts plain text from a Word/Excel/PowerPoint file for classification."""
    n = name.lower()
    try:
        if n.endswith(".docx"):
            from docx import Document
            return "\n".join(p.text for p in Document(io.BytesIO(data)).paragraphs)[:65536]
        if n.endswith(".xlsx"):
            from openpyxl import load_workbook
            wb = load_workbook(io.BytesIO(data), read_only=True, data_only=True)
            parts = []
            for ws in wb.worksheets:
                for row in ws.iter_rows(values_only=True):
                    parts.append(" ".join(str(c) for c in row if c is not None))
                    if sum(len(p) for p in parts) > 65536:
                        break
            return " ".join(parts)[:65536]
        if n.endswith(".pptx"):
            from pptx import Presentation
            out = []
            for slide in Presentation(io.BytesIO(data)).slides:
                for s in slide.shapes:
                    if getattr(s, "has_text_frame", False):
                        out.append(s.text_frame.text)
            return "\n".join(out)[:65536]
    except Exception:
        return ""
    return ""


def scan_tenant(db: Session, cfg: models.MicrosoftConfig) -> dict:
    if not (cfg.azure_tenant_id and cfg.client_id and cfg.client_secret and cfg.drive_user):
        return _finish(db, cfg, "error: set Azure tenant id, client id/secret and a drive user")
    try:
        token = get_token(cfg.azure_tenant_id, cfg.client_id, cfg.client_secret)
    except Exception as e:
        return _finish(db, cfg, f"error: auth failed ({str(e)[:120]})")

    pol = db.query(models.StampPolicy).filter(models.StampPolicy.tenant_id == cfg.tenant_id).first()
    template = pol.text_template if pol else "CLASSIFICATION: {label}"
    placement = cfg.placement or "footer"
    done = {s.file_id for s in db.query(models.StampedDoc.file_id)
            .filter(models.StampedDoc.tenant_id == cfg.tenant_id,
                    models.StampedDoc.provider == "microsoft").all()}

    try:
        docs = list_office_docs(token, cfg.drive_user)
    except Exception as e:
        return _finish(db, cfg, f"error: listing OneDrive failed ({str(e)[:120]})")

    scanned = stamped = 0
    for f in docs:
        if f["id"] in done:
            continue
        scanned += 1
        try:
            data = download(token, cfg.drive_user, f["id"])
        except Exception:
            continue
        name = f.get("name", "")
        text = doc_text(data, name)
        label, _ = classify_text(db, cfg.tenant_id, name, text)
        if not label:
            continue
        try:
            new = stamp_office(data, name, template.replace("{label}", label.name), label.color, placement)
            upload(token, cfg.drive_user, f["id"], new)
            db.add(models.StampedDoc(tenant_id=cfg.tenant_id, provider="microsoft",
                                     file_id=f["id"], label=label.name))
            stamped += 1
        except Exception:
            continue
    db.commit()
    return _finish(db, cfg, f"ok: scanned {scanned} new files, stamped {stamped}")


def _finish(db: Session, cfg: models.MicrosoftConfig, status: str) -> dict:
    cfg.last_scan = datetime.now(timezone.utc)
    cfg.last_status = status[:255]
    db.commit()
    return {"status": status}
